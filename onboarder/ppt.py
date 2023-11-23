from pptx import Presentation


def extract_text_from_ppt(ppt_path):
    # PowerPointのプレゼンテーションを開く
    try:
        prs = Presentation(ppt_path)
        # スライドごとにテキストを抽出する
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        return text
    except FileNotFoundError:
        print("PowerPointファイルが見つかりませんでした")
    except ValueError:
        print("PowerPointファイルを読み込めませんでした")
    return ""
